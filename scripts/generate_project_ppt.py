#!/usr/bin/env python3
"""Generate the AI Werewolf project report deck.

This script intentionally keeps the deck content close to the acceptance
documents under docs/. It generates a PowerPoint file and a lightweight
Markdown outline for review.
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable
from typing import Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "presentations"
ASSET_DIR = ROOT / "docs" / "assets" / "final_report" / "png"
PPTX_PATH = OUT_DIR / "AI_Werewolf_Project_Report.pptx"
OUTLINE_PATH = OUT_DIR / "AI_Werewolf_Project_Report_outline.md"

FONT = "Microsoft YaHei"
FONT_LATIN = "Aptos"


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


COLORS = {
    "bg": rgb("#101722"),
    "bg2": rgb("#142033"),
    "panel": rgb("#182538"),
    "panel2": rgb("#203047"),
    "ink": rgb("#F7F2E8"),
    "muted": rgb("#B8C2D1"),
    "line": rgb("#38516F"),
    "gold": rgb("#D9A441"),
    "red": rgb("#C8524A"),
    "green": rgb("#5FB878"),
    "blue": rgb("#5CA7D8"),
    "purple": rgb("#9B7BD8"),
    "white": rgb("#FFFFFF"),
}


def emu(value: float) -> int:
    return Inches(value)


def set_background(slide, color: RGBColor = COLORS["bg"]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    color: RGBColor = COLORS["ink"],
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font: str = FONT,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = emu(margin)
    tf.margin_right = emu(margin)
    tf.margin_top = emu(margin)
    tf.margin_bottom = emu(margin)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_bullets(
    slide,
    items: Sequence[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 15,
    color: RGBColor = COLORS["ink"],
    bullet_color: RGBColor = COLORS["gold"],
    spacing: int = 7,
):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.06)
    tf.margin_right = emu(0.04)
    tf.margin_top = emu(0.04)
    tf.margin_bottom = emu(0.04)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        p.space_after = Pt(spacing)
        dot = p.add_run()
        dot.text = "• "
        dot.font.name = FONT
        dot.font.size = Pt(size)
        dot.font.bold = True
        dot.font.color.rgb = bullet_color
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None, section: str | None = None) -> None:
    add_text(slide, title, 0.62, 0.35, 8.3, 0.55, size=25, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.65, 0.92, 8.8, 0.28, size=10, color=COLORS["muted"])
    if section:
        add_pill(slide, section, 10.75, 0.43, 1.85, 0.36, fill=COLORS["panel2"], line=COLORS["line"], size=10)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(0.62), emu(1.22), emu(1.2), emu(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["gold"]
    line.line.fill.background()


def add_footer(slide, page: int) -> None:
    add_text(slide, "AI Werewolf · Play / Evaluate / Evolve", 0.62, 7.06, 3.8, 0.22, size=8, color=COLORS["muted"])
    add_text(slide, f"{page:02d}", 12.05, 7.04, 0.55, 0.25, size=9, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = COLORS["panel"],
    line: RGBColor = COLORS["line"],
    radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    shape = slide.shapes.add_shape(radius_shape, emu(x), emu(y), emu(w), emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = COLORS["panel"],
    line: RGBColor = COLORS["line"],
    color: RGBColor = COLORS["ink"],
    size: int = 11,
    bold: bool = False,
):
    shape = add_card(slide, x, y, w, h, fill=fill, line=line)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = emu(0.05)
    tf.margin_right = emu(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return shape


def add_metric(
    slide,
    label: str,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    accent: RGBColor = COLORS["gold"],
    note: str | None = None,
):
    add_card(slide, x, y, w, h, fill=COLORS["panel"], line=COLORS["line"])
    add_text(slide, value, x + 0.16, y + 0.14, w - 0.32, h * 0.38, size=25, bold=True, color=accent)
    add_text(slide, label, x + 0.18, y + 0.78, w - 0.36, 0.28, size=10, color=COLORS["ink"], bold=True)
    if note:
        add_text(slide, note, x + 0.18, y + 1.1, w - 0.36, h - 1.15, size=8, color=COLORS["muted"])


def add_image_fit(slide, image: Path, x: float, y: float, w: float, h: float):
    with Image.open(image) as img:
        iw, ih = img.size
    image_ratio = iw / ih
    box_ratio = w / h
    if image_ratio >= box_ratio:
        draw_w = w
        draw_h = w / image_ratio
        dx = 0
        dy = (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * image_ratio
        dx = (w - draw_w) / 2
        dy = 0
    return slide.shapes.add_picture(str(image), emu(x + dx), emu(y + dy), width=emu(draw_w), height=emu(draw_h))


def add_table(
    slide,
    rows: Sequence[Sequence[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    header_fill: RGBColor = COLORS["panel2"],
    body_fill: RGBColor = COLORS["panel"],
    font_size: int = 9,
    widths: Sequence[float] | None = None,
):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), emu(x), emu(y), emu(w), emu(h))
    table = shape.table
    if widths:
        for idx, width in enumerate(widths):
            table.columns[idx].width = emu(width)
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r_idx == 0 else body_fill
            cell.margin_left = emu(0.04)
            cell.margin_right = emu(0.04)
            cell.margin_top = emu(0.03)
            cell.margin_bottom = emu(0.03)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(font_size if r_idx else font_size + 1)
                p.font.bold = r_idx == 0
                p.font.color.rgb = COLORS["ink"] if r_idx == 0 else COLORS["muted"]
                p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
    return shape


def add_flow(slide, labels: Sequence[str], x: float, y: float, w: float, *, color: RGBColor = COLORS["blue"]) -> None:
    box_w = (w - 0.9) / len(labels)
    for idx, label in enumerate(labels):
        bx = x + idx * (box_w + 0.3)
        add_pill(
            slide, label, bx, y, box_w, 0.55, fill=COLORS["panel"], line=color, color=COLORS["ink"], size=11, bold=True
        )
        if idx < len(labels) - 1:
            add_text(slide, "→", bx + box_w + 0.07, y + 0.1, 0.16, 0.22, size=18, color=color, bold=True)


def numbered_section(slide, number: str, title: str, body: str, x: float, y: float, w: float, accent: RGBColor) -> None:
    add_card(slide, x, y, w, 1.08, fill=COLORS["panel"], line=COLORS["line"])
    add_text(slide, number, x + 0.18, y + 0.18, 0.45, 0.3, size=16, bold=True, color=accent)
    add_text(slide, title, x + 0.68, y + 0.16, w - 0.85, 0.25, size=14, bold=True)
    add_text(slide, body, x + 0.68, y + 0.49, w - 0.85, 0.32, size=9, color=COLORS["muted"])


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, "AI Werewolf", 0.72, 0.62, 5.4, 0.55, size=36, bold=True)
    add_text(slide, "多智能体狼人杀项目汇报", 0.75, 1.26, 4.2, 0.38, size=18, color=COLORS["gold"], bold=True)
    add_text(
        slide,
        "对战 Play · 复盘 Evaluate · 进化 Evolve",
        0.76,
        1.78,
        5.8,
        0.34,
        size=16,
        color=COLORS["muted"],
    )
    add_bullets(
        slide,
        [
            "严格信息隔离下的多 Agent 博弈系统",
            "完整对局引擎 + 观战 UI + 真人/AI 混战入口",
            "Track B 决策复盘评价，Track C 策略知识回流",
        ],
        0.8,
        2.55,
        4.8,
        1.3,
        size=14,
    )
    add_pill(slide, "2026-06-08", 0.82, 6.36, 1.35, 0.34, fill=COLORS["panel2"], line=COLORS["line"], size=9)
    add_pill(slide, "wxhfy / 付一涵", 2.3, 6.36, 1.6, 0.34, fill=COLORS["panel2"], line=COLORS["line"], size=9)
    add_image_fit(slide, ASSET_DIR / "play-evaluate-evolve.png", 6.15, 0.55, 6.45, 5.75)
    add_footer(slide, 1)


def slide_agenda(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "汇报主线", "从产品定位、系统架构、核心创新到验收证据", "Overview")
    numbered_section(
        slide, "01", "项目定位", "为什么狼人杀适合作为信息不对称多智能体研究平台", 0.8, 1.75, 5.55, COLORS["gold"]
    )
    numbered_section(
        slide, "02", "系统架构", "前后端实时系统、规则引擎、Agent 决策与数据库证据链", 6.8, 1.75, 5.55, COLORS["blue"]
    )
    numbered_section(
        slide,
        "03",
        "核心创新",
        "信息隔离、CognitiveAgent、Track B 复盘、Track C 自进化",
        0.8,
        3.25,
        5.55,
        COLORS["green"],
    )
    numbered_section(
        slide, "04", "验收与边界", "自动化测试、端到端 smoke、风险项与下一步计划", 6.8, 3.25, 5.55, COLORS["red"]
    )
    add_card(slide, 0.8, 5.38, 11.55, 0.68, fill=COLORS["bg2"], line=COLORS["line"])
    add_text(
        slide,
        "一句话结论：项目已形成可演示、可复盘、可积累策略知识的 AI 狼人杀闭环系统。",
        1.08,
        5.55,
        10.8,
        0.25,
        size=14,
        bold=True,
        color=COLORS["ink"],
    )
    add_footer(slide, 2)


def slide_positioning(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "项目定位", "AI Werewolf = 多智能体狼人杀研究平台", "Positioning")
    add_text(
        slide,
        "研究问题",
        0.84,
        1.58,
        1.8,
        0.28,
        size=15,
        bold=True,
        color=COLORS["gold"],
    )
    add_text(
        slide,
        "在信息不对称、阵营对抗和自然语言交流同时存在的环境里，如何让 Agent 既遵守规则，又具备角色化推理、社交判断、复盘学习能力？",
        0.84,
        1.98,
        5.25,
        1.0,
        size=17,
        bold=True,
    )
    features = [
        ("完整对局", "15+ 细分阶段，夜晚技能、白天发言、投票、胜负判定"),
        ("严格隔离", "GameState 与 PlayerView 分离，Agent 只看到身份允许的信息"),
        ("角色化 Agent", "MBTI 人格 + Role 身份 + Strategy 知识三层架构"),
        ("可复盘进化", "每步决策可审计、可解释，并沉淀为下一局策略"),
    ]
    y = 3.48
    for idx, (title, body) in enumerate(features):
        x = 0.84 + (idx % 2) * 3.0
        yy = y + (idx // 2) * 1.15
        add_card(slide, x, yy, 2.65, 0.88, fill=COLORS["panel"], line=COLORS["line"])
        add_text(slide, title, x + 0.16, yy + 0.13, 2.35, 0.22, size=12, bold=True, color=COLORS["gold"])
        add_text(slide, body, x + 0.16, yy + 0.43, 2.35, 0.28, size=8, color=COLORS["muted"])
    add_card(slide, 6.75, 1.55, 5.7, 4.85, fill=COLORS["panel"], line=COLORS["line"])
    add_text(slide, "两条主线", 7.1, 1.85, 1.5, 0.3, size=15, bold=True, color=COLORS["gold"])
    add_flow(slide, ["Play", "Evaluate", "Evolve"], 7.1, 2.45, 4.95, color=COLORS["green"])
    add_text(
        slide,
        "系统闭环：对局执行 → 赛后复盘 → 经验抽取 → 下一局策略回流",
        7.12,
        3.22,
        4.85,
        0.34,
        size=10,
        color=COLORS["muted"],
    )
    add_flow(slide, ["Persona", "Role", "Strategy"], 7.1, 4.08, 4.95, color=COLORS["blue"])
    add_text(
        slide,
        "Agent 认知层：人格决定风格，身份决定规则边界，策略决定如何赢",
        7.12,
        4.86,
        4.85,
        0.34,
        size=10,
        color=COLORS["muted"],
    )
    add_footer(slide, 3)


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "系统总体架构", "前端实时观战、后端规则引擎、Agent 决策、B/C 评测进化闭环", "Architecture")
    add_image_fit(slide, ASSET_DIR / "system-architecture.png", 0.55, 1.5, 8.15, 5.2)
    add_text(slide, "架构要点", 9.05, 1.55, 1.5, 0.3, size=15, bold=True, color=COLORS["gold"])
    add_bullets(
        slide,
        [
            "FastAPI + WebSocket 承担房间、快照和实时事件推送",
            "WerewolfGame 是唯一规则主控，Agent 只提交 Decision",
            "PostgreSQL 保存事件、快照、决策、复盘和知识",
            "Track B/C 在赛后消费证据链，再回流到 StrategyRetriever",
        ],
        9.0,
        2.05,
        3.55,
        2.05,
        size=11,
    )
    add_metric(
        slide,
        "核心 ORM / 证据表",
        "20+",
        9.05,
        4.65,
        1.6,
        1.3,
        accent=COLORS["blue"],
        note="games / events / decisions / reviews / knowledge",
    )
    add_metric(
        slide,
        "主要页面路由",
        "7",
        10.88,
        4.65,
        1.6,
        1.3,
        accent=COLORS["green"],
        note="大厅、对局、Human、复盘、看板、Persona",
    )
    add_footer(slide, 4)


def slide_game_flow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "单局对局流程", "夜晚技能、白天发言投票、胜负判定与赛后处理", "Game Flow")
    add_image_fit(slide, ASSET_DIR / "game-operation-flow.png", 0.65, 1.45, 7.05, 5.15)
    add_text(slide, "阶段流转", 8.0, 1.55, 1.6, 0.28, size=15, bold=True, color=COLORS["gold"])
    rows = [
        ["阶段", "关键动作", "主控"],
        ["夜晚", "守卫 / 狼刀 / 女巫 / 预言家", "引擎结算"],
        ["白天", "发言 / 警长 / 投票 / 遗言", "阶段机调度"],
        ["终局", "胜负判定 / 复盘 / 知识抽取", "B/C 后处理"],
    ]
    add_table(slide, rows, 8.0, 2.03, 4.45, 1.75, font_size=8, widths=[1.0, 2.2, 1.25])
    add_bullets(
        slide,
        [
            "Agent 不直接修改 GameState，所有行动先过合法性校验",
            "夜晚并行动作在结算阶段统一处理，避免先后执行悖论",
            "终局后自动形成可回放、可审计、可抽取知识的证据链",
        ],
        8.05,
        4.15,
        4.25,
        1.45,
        size=11,
    )
    add_footer(slide, 5)


def slide_modules(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "核心模块清单", "从引擎到前端的可交付模块", "Modules")
    add_image_fit(slide, ASSET_DIR / "module-map.png", 0.55, 1.42, 6.7, 5.15)
    rows = [
        ["模块", "职责", "状态"],
        ["WerewolfGame", "初始化、阶段推进、结算、胜负、事件记录", "通过"],
        ["Visibility", "GameState -> PlayerView，信息裁剪", "92/92"],
        ["CognitiveAgent", "Observe / Memory / Planner / Tools / LLM", "通过"],
        ["Track B", "决策审计、复盘、证据引用", "通过"],
        ["Track C", "lesson 抽取、知识生命周期、检索回流", "通过"],
        ["Frontend", "大厅、观战、Human、复盘、仪表盘", "通过"],
    ]
    add_table(slide, rows, 7.55, 1.55, 5.05, 4.25, font_size=7, widths=[1.35, 2.9, 0.8])
    add_text(
        slide,
        "设计原则：规则由引擎控制，视图由 Visibility 裁剪，Agent 输出意图，所有行为进入审计链。",
        7.62,
        6.05,
        4.9,
        0.36,
        size=10,
        color=COLORS["muted"],
    )
    add_footer(slide, 6)


def slide_visibility(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(
        slide, "信息隔离：可信博弈的底座", "GameState 是真相，PlayerView 是每个 Agent 的合法局部视角", "Visibility"
    )
    add_metric(
        slide, "边界检查", "92/92", 0.75, 1.55, 2.15, 1.35, accent=COLORS["green"], note="verify_visibility_strict.py"
    )
    add_metric(
        slide, "公开角色泄漏", "0", 3.15, 1.55, 2.15, 1.35, accent=COLORS["green"], note="demo public role leak=false"
    )
    add_metric(slide, "视图体系", "3", 5.55, 1.55, 2.15, 1.35, accent=COLORS["blue"], note="观众 / 上帝 / 玩家")
    add_card(slide, 0.75, 3.35, 6.95, 2.75, fill=COLORS["panel"], line=COLORS["line"])
    add_text(slide, "隔离边界", 1.05, 3.65, 1.5, 0.28, size=14, bold=True, color=COLORS["gold"])
    add_bullets(
        slide,
        [
            "村民不可见查验、女巫夜间受害者、狼队沟通",
            "狼人只可见狼队队友，不可见神职私有结果",
            "预言家只可见自己的查验，女巫只可见自己的药与受害者",
            "public event 不泄露角色分配或查验结果",
        ],
        1.0,
        4.12,
        6.25,
        1.25,
        size=11,
    )
    add_card(slide, 8.25, 1.58, 4.1, 4.52, fill=COLORS["panel"], line=COLORS["line"])
    add_text(slide, "GameState", 9.55, 1.95, 1.55, 0.3, size=15, bold=True, color=COLORS["red"], align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "↓ Visibility.for_player(player_id)",
        8.85,
        2.55,
        2.9,
        0.25,
        size=11,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_pill(slide, "PlayerView P1", 8.65, 3.2, 1.2, 0.48, fill=COLORS["bg2"], line=COLORS["blue"], size=9)
    add_pill(slide, "PlayerView P2", 10.0, 3.2, 1.2, 0.48, fill=COLORS["bg2"], line=COLORS["green"], size=9)
    add_pill(slide, "PlayerView P3", 9.32, 4.0, 1.2, 0.48, fill=COLORS["bg2"], line=COLORS["gold"], size=9)
    add_text(
        slide,
        "每个 Agent 的输入都由代码裁剪，而不是靠 Prompt 约束“不要偷看”。",
        8.65,
        4.95,
        3.25,
        0.55,
        size=10,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 7)


def slide_agent(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(
        slide, "CognitiveAgent：角色化决策主体", "Observe -> Think -> Act -> Reflect，工具调用 trace 可审计", "Agent"
    )
    add_flow(slide, ["Observe", "Think", "Act", "Reflect"], 0.85, 1.65, 6.2, color=COLORS["gold"])
    add_text(slide, "三层 Prompt / 认知上下文", 0.88, 2.65, 2.6, 0.3, size=14, bold=True, color=COLORS["gold"])
    rows = [
        ["层", "作用"],
        ["Persona / MBTI", "控制表达风格、认知倾向、风险偏好"],
        ["Role Identity", "定义身份、阵营、技能和胜利条件"],
        ["Strategy + Tools", "注入当前可用策略、反模式和检索结果"],
    ]
    add_table(slide, rows, 0.85, 3.1, 5.45, 1.75, font_size=8, widths=[1.45, 4.0])
    add_text(slide, "工具调用循环", 7.0, 1.62, 1.8, 0.3, size=14, bold=True, color=COLORS["gold"])
    tools = [
        "search_strategies",
        "recall_memory",
        "check_rules",
        "get_social_info",
        "analyze_votes",
        "set_strategic_intent",
        "submit_decision",
    ]
    for idx, tool in enumerate(tools):
        x = 7.05 + (idx % 2) * 2.45
        y = 2.12 + (idx // 2) * 0.7
        add_pill(slide, tool, x, y, 2.15, 0.38, fill=COLORS["panel"], line=COLORS["line"], size=8)
    add_card(slide, 7.05, 5.15, 4.9, 0.72, fill=COLORS["bg2"], line=COLORS["line"])
    add_text(
        slide,
        "收益：角色差异、记忆连续性、社交判断和策略回流都能进入同一条可审计决策链。",
        7.25,
        5.35,
        4.5,
        0.24,
        size=10,
        color=COLORS["ink"],
        bold=True,
    )
    add_footer(slide, 8)


def slide_loop(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Play → Evaluate → Evolve", "系统不只跑一局，还能解释一局、积累一局", "Closed Loop")
    add_image_fit(slide, ASSET_DIR / "play-evaluate-evolve.png", 0.7, 1.45, 6.95, 4.8)
    add_card(slide, 8.05, 1.65, 4.15, 1.2, fill=COLORS["panel"], line=COLORS["line"])
    add_text(slide, "Play", 8.32, 1.86, 1.0, 0.26, size=16, bold=True, color=COLORS["green"])
    add_text(
        slide,
        "完整对局执行，写入 GameEvent / Snapshot / AgentDecision",
        9.25,
        1.83,
        2.55,
        0.36,
        size=10,
        color=COLORS["muted"],
    )
    add_card(slide, 8.05, 3.12, 4.15, 1.2, fill=COLORS["panel"], line=COLORS["line"])
    add_text(slide, "Evaluate", 8.32, 3.33, 1.25, 0.26, size=16, bold=True, color=COLORS["blue"])
    add_text(
        slide,
        "Track B 评价关键决策，生成 ScoredStep 和 PublishedReview",
        9.55,
        3.3,
        2.25,
        0.36,
        size=10,
        color=COLORS["muted"],
    )
    add_card(slide, 8.05, 4.6, 4.15, 1.2, fill=COLORS["panel"], line=COLORS["line"])
    add_text(slide, "Evolve", 8.32, 4.81, 1.25, 0.26, size=16, bold=True, color=COLORS["gold"])
    add_text(
        slide,
        "Track C 抽取 lesson，进入 candidate / active 知识池再被检索",
        9.45,
        4.78,
        2.35,
        0.36,
        size=10,
        color=COLORS["muted"],
    )
    add_footer(slide, 9)


def slide_track_b(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Track B：决策复盘评价", "从“谁赢了”推进到“每一步为什么好/坏”", "Evaluate")
    add_text(slide, "三级评价链路", 0.78, 1.55, 2.0, 0.3, size=15, bold=True, color=COLORS["gold"])
    add_flow(slide, ["Tier 1 规则", "Tier 2 轻量 LLM", "Tier 3 复核"], 0.78, 2.08, 6.4, color=COLORS["blue"])
    add_bullets(
        slide,
        [
            "明确场景用规则：投票、技能目标、信息泄露、非法行动",
            "模糊场景用轻量 LLM：发言有效性、推理质量、局势影响",
            "高影响决策可进入复核流程，保留 counterfactual 和证据引用",
        ],
        0.82,
        3.05,
        6.2,
        1.45,
        size=12,
    )
    add_text(slide, "输出对象", 0.82, 5.25, 1.5, 0.25, size=14, bold=True, color=COLORS["gold"])
    add_flow(
        slide, ["DecisionScore", "ScoredStep", "PlayerReview", "PublishedReview"], 2.0, 5.15, 5.4, color=COLORS["green"]
    )
    add_metric(
        slide,
        "B/C 验收专项",
        "42 passed",
        8.0,
        1.7,
        2.0,
        1.35,
        accent=COLORS["green"],
        note="test_api / test_b_full_acceptance / test_track_c",
    )
    add_metric(
        slide,
        "strict ScoredStep 覆盖",
        "27/27",
        10.28,
        1.7,
        2.0,
        1.35,
        accent=COLORS["blue"],
        note="backend_acceptance 记录",
    )
    add_card(slide, 8.0, 3.65, 4.3, 1.78, fill=COLORS["panel"], line=COLORS["line"])
    add_text(slide, "可展示价值", 8.28, 3.95, 1.4, 0.24, size=13, bold=True, color=COLORS["gold"])
    add_text(
        slide,
        "复盘报告不只是文本总结，而是可回链到原始事件、Agent 决策和当时可见信息的证据化报告。",
        8.28,
        4.35,
        3.65,
        0.55,
        size=11,
        color=COLORS["muted"],
    )
    add_footer(slide, 10)


def slide_track_c(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Track C：知识抽取与安全回流", "把复盘结论沉淀为下一局可检索策略", "Evolve")
    add_image_fit(slide, ASSET_DIR / "track-bc-flow.png", 0.62, 1.45, 6.95, 4.85)
    add_text(slide, "知识生命周期", 8.0, 1.6, 2.0, 0.3, size=15, bold=True, color=COLORS["gold"])
    add_flow(slide, ["candidate", "active", "deprecated"], 8.0, 2.15, 4.25, color=COLORS["gold"])
    add_bullets(
        slide,
        [
            "新 lesson 默认进入 candidate，避免直接污染 active 策略池",
            "检索侧使用 confidence / visibility / privacy / applicability 四重过滤",
            "保留 source_game / source_event / source_item，知识可追溯",
            "Track C 只注入 Strategy 层，不污染身份规则和人格风格",
        ],
        8.02,
        3.1,
        4.25,
        1.65,
        size=11,
    )
    add_metric(slide, "当前 active", "401", 8.05, 5.35, 1.3, 1.1, accent=COLORS["green"], note="DB 快照")
    add_metric(slide, "当前 candidate", "3856", 9.55, 5.35, 1.3, 1.1, accent=COLORS["gold"], note="DB 快照")
    add_metric(slide, "usage feedback", "127k", 11.05, 5.35, 1.3, 1.1, accent=COLORS["blue"], note="DB 快照")
    add_footer(slide, 11)


def slide_frontend(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "产品化界面", "Next.js 观战 UI、Human 模式、复盘与评测仪表盘", "Frontend")
    pages = [
        ("/", "大厅 / 创建房间", COLORS["gold"]),
        ("/room/[id]/play", "AI 对局观战", COLORS["blue"]),
        ("/room/[id]/human", "真人玩家操作", COLORS["green"]),
        ("/games/[id]/report", "单局复盘报告", COLORS["purple"]),
        ("/eval/dashboard", "评测仪表盘", COLORS["blue"]),
        ("/personas", "Persona 库", COLORS["gold"]),
    ]
    for idx, (route, label, accent) in enumerate(pages):
        x = 0.82 + (idx % 2) * 5.85
        y = 1.55 + (idx // 2) * 1.05
        add_card(slide, x, y, 5.2, 0.72, fill=COLORS["panel"], line=accent)
        add_text(slide, route, x + 0.18, y + 0.12, 1.95, 0.24, size=10, color=accent, bold=True, font=FONT_LATIN)
        add_text(slide, label, x + 2.2, y + 0.12, 2.65, 0.24, size=12, color=COLORS["ink"], bold=True)
        add_text(slide, "已纳入 build / UI smoke 验证", x + 2.2, y + 0.4, 2.65, 0.18, size=8, color=COLORS["muted"])
    add_card(slide, 0.82, 5.95, 11.05, 0.5, fill=COLORS["bg2"], line=COLORS["line"])
    add_text(
        slide,
        "UI smoke 覆盖：对局到 GAME_END、复盘 iframe、HTML 复盘资产、中英文切换、AI/Human 房间流程。",
        1.05,
        6.1,
        10.65,
        0.2,
        size=10,
        color=COLORS["ink"],
    )
    add_footer(slide, 12)


def slide_acceptance(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "自动化验收结果", "当前工作区本地可复现验证，fake LLM provider，日期 2026-06-08", "Acceptance")
    metrics = [
        ("pytest tests/", "449 passed", "27 skipped, 26 warnings"),
        ("E2E smoke", "passed", "scripts/e2e_smoke.py"),
        ("Visibility strict", "92 / 0", "92 passed, 0 failed"),
        ("Demo game", "GAME_END", "10 players, 157 events"),
        ("Frontend", "lint + build", "UI smoke passed"),
        ("Track B/C", "42 passed", "关键链路专项"),
    ]
    for idx, (label, value, note) in enumerate(metrics):
        x = 0.78 + (idx % 3) * 4.05
        y = 1.55 + (idx // 3) * 1.65
        accent = [COLORS["green"], COLORS["blue"], COLORS["gold"], COLORS["purple"], COLORS["green"], COLORS["blue"]][
            idx
        ]
        add_metric(slide, label, value, x, y, 3.35, 1.3, accent=accent, note=note)
    rows = [
        ["维度", "结论", "证据"],
        ["对局引擎", "通过", "backend.run_demo 到 GAME_END"],
        ["信息隔离", "通过", "92 项检查，public role leak=false"],
        ["后端 API", "通过", "rooms / games / snapshot / health smoke"],
        ["前端流程", "通过", "build + browser smoke"],
        ["复盘与进化", "通过", "报告导出、B->C、知识回流"],
    ]
    add_table(slide, rows, 0.78, 5.05, 11.55, 1.45, font_size=7, widths=[1.75, 1.0, 8.8])
    add_footer(slide, 13)


def slide_fixes_risks(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "本轮修复与风险边界", "验收中发现的问题已收敛，同时保留真实 LLM 长跑等补验项", "Risks")
    add_text(slide, "已修复", 0.82, 1.55, 1.2, 0.3, size=15, bold=True, color=COLORS["green"])
    fixed = [
        "恢复 Track B speech semantic audit 与 speech act classifier 报告",
        "补齐 StrategyKnowledgeDocData experiment_id 传递",
        "修复 AbstractedLesson 到 PostgreSQL 字段映射",
        "修复知识 upsert 时间序列化与 UI smoke 断言",
    ]
    add_bullets(slide, fixed, 0.82, 2.0, 5.65, 1.8, size=11, bullet_color=COLORS["green"])
    add_text(slide, "风险与补验", 7.0, 1.55, 1.6, 0.3, size=15, bold=True, color=COLORS["gold"])
    risks = [
        "答辩前建议用真实 provider 跑 7P / 9P / 12P 各 1 局",
        "Track C 胜率提升仍需 paired seed 在线 A/B 实验",
        "并发 WebSocket / 重连压力测试未作为本轮自动化主项",
        "数据库历史实验需按 experiment_id / provider / strict 标记过滤",
    ]
    add_bullets(slide, risks, 7.0, 2.0, 5.15, 1.8, size=11, bullet_color=COLORS["gold"])
    add_card(slide, 0.82, 4.65, 11.3, 1.05, fill=COLORS["bg2"], line=COLORS["line"])
    add_text(slide, "结论边界", 1.1, 4.92, 1.3, 0.25, size=13, bold=True, color=COLORS["gold"])
    add_text(
        slide,
        "当前 PPT 中的验收结论引用的是本地可复现自动化结果；历史数据库快照和占位图表不混写为正式提升结论。",
        2.25,
        4.9,
        9.55,
        0.28,
        size=10,
        color=COLORS["ink"],
    )
    add_footer(slide, 14)


def slide_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "总结与下一步", "从可玩系统到可验证研究平台", "Summary")
    add_text(slide, "当前交付", 0.85, 1.6, 1.6, 0.3, size=15, bold=True, color=COLORS["gold"])
    deliverables = [
        ("能玩", "完整狼人杀引擎，支持 AI 对局、Human 模式、实时观战"),
        ("能解释", "事件、决策、复盘全链路可追溯"),
        ("能积累", "Track C 将复盘经验转为可检索策略知识"),
    ]
    for idx, (title, body) in enumerate(deliverables):
        y = 2.1 + idx * 1.08
        add_card(
            slide,
            0.88,
            y,
            5.15,
            0.78,
            fill=COLORS["panel"],
            line=[COLORS["green"], COLORS["blue"], COLORS["gold"]][idx],
        )
        add_text(
            slide,
            title,
            1.15,
            y + 0.17,
            0.7,
            0.22,
            size=14,
            bold=True,
            color=[COLORS["green"], COLORS["blue"], COLORS["gold"]][idx],
        )
        add_text(slide, body, 2.0, y + 0.18, 3.65, 0.22, size=10, color=COLORS["muted"])
    add_text(slide, "下一步", 7.0, 1.6, 1.6, 0.3, size=15, bold=True, color=COLORS["gold"])
    next_steps = [
        "冻结一次真实 LLM + PostgreSQL strict 验收输出",
        "补 paired seed 多局 A/B，量化 Track C 对胜率/决策分的影响",
        "扩展并发观战、断线重连和更多角色规则专项测试",
        "将复盘报告与策略卡进一步产品化，形成演示闭环",
    ]
    add_bullets(slide, next_steps, 7.0, 2.1, 4.95, 2.05, size=11, bullet_color=COLORS["gold"])
    add_card(slide, 0.88, 5.65, 11.25, 0.72, fill=COLORS["bg2"], line=COLORS["line"])
    add_text(
        slide,
        "最终判断：当前项目主体功能达到交付演示与答辩验收要求。",
        1.18,
        5.86,
        10.65,
        0.24,
        size=15,
        bold=True,
        color=COLORS["ink"],
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 15)


def build_outline() -> str:
    slides = [
        ("封面", "AI Werewolf 多智能体狼人杀项目汇报"),
        ("汇报主线", "项目定位、系统架构、核心创新、验收边界"),
        ("项目定位", "对战 Play、复盘 Evaluate、进化 Evolve 的研究平台"),
        ("系统总体架构", "前端、FastAPI/WebSocket、WerewolfGame、Agent、DB、Track B/C"),
        ("单局对局流程", "夜晚、白天、终局与赛后处理"),
        ("核心模块清单", "引擎、Visibility、Agent、Track B、Track C、前端"),
        ("信息隔离", "GameState 与 PlayerView 分离，92/92 边界检查"),
        ("CognitiveAgent", "三层 Prompt 与工具调用式决策循环"),
        ("闭环主线", "Play -> Evaluate -> Evolve -> Retrieve"),
        ("Track B", "决策复盘评价、报告生成、证据引用"),
        ("Track C", "知识生命周期、安全回流与 Wiki/Hermes 增量设计"),
        ("产品化界面", "大厅、对局、Human、复盘、仪表盘、Persona"),
        ("自动化验收", "pytest、E2E、visibility、frontend、demo、B/C 专项"),
        ("修复与风险", "本轮修复项、真实 LLM 长跑和 A/B 补验边界"),
        ("总结与下一步", "能玩、能解释、能积累；补真实验收与多局 A/B"),
    ]
    lines = [
        "# AI Werewolf 项目汇报 PPT 大纲",
        "",
        "> 生成文件：`AI_Werewolf_Project_Report.pptx`",
        "> 生成日期：2026-06-08",
        "",
    ]
    for idx, (title, body) in enumerate(slides, 1):
        lines.append(f"## {idx}. {title}")
        lines.append("")
        lines.append(body)
        lines.append("")
    return "\n".join(lines)


def build_deck() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = emu(13.333)
    prs.slide_height = emu(7.5)

    builders: Iterable = [
        slide_cover,
        slide_agenda,
        slide_positioning,
        slide_architecture,
        slide_game_flow,
        slide_modules,
        slide_visibility,
        slide_agent,
        slide_loop,
        slide_track_b,
        slide_track_c,
        slide_frontend,
        slide_acceptance,
        slide_fixes_risks,
        slide_summary,
    ]
    for builder in builders:
        builder(prs)

    prs.save(PPTX_PATH)
    OUTLINE_PATH.write_text(build_outline(), encoding="utf-8")


if __name__ == "__main__":
    build_deck()
    print(PPTX_PATH)
    print(OUTLINE_PATH)
